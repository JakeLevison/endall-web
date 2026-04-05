#!/usr/bin/env python3
"""Generate static demo files for /public/demo-files/.

Outputs realistic-looking workbooks, proposals, and reports for
Patriot Electric (Ashburn, VA - data center subcontractor).
"""

from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from fpdf import FPDF

OUT = Path(__file__).resolve().parents[1] / "public" / "demo-files"
OUT.mkdir(parents=True, exist_ok=True)

# --- Styling helpers -------------------------------------------------------
HDR_FILL = PatternFill(start_color="1F2937", end_color="1F2937", fill_type="solid")
HDR_FONT = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
SECTION_FILL = PatternFill(start_color="E5E7EB", end_color="E5E7EB", fill_type="solid")
THIN = Side(style="thin", color="CBD5E1")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)


def style_header_row(ws, row, cols):
    for col in range(1, cols + 1):
        c = ws.cell(row=row, column=col)
        c.fill = HDR_FILL
        c.font = HDR_FONT
        c.alignment = Alignment(horizontal="left", vertical="center")
        c.border = BORDER


def auto_width(ws, cols):
    for i, col in enumerate(cols, start=1):
        letter = ws.cell(row=1, column=i).column_letter
        max_len = max([len(str(col))] + [len(str(ws.cell(row=r, column=i).value or "")) for r in range(2, ws.max_row + 1)])
        ws.column_dimensions[letter].width = max(12, min(40, max_len + 2))


# --- 1. NPV Analysis -------------------------------------------------------
def build_npv():
    wb = Workbook()
    ws = wb.active
    ws.title = "How to Use"
    ws["A1"] = "Patriot Electric - NPV Analysis"
    ws["A1"].font = Font(bold=True, size=16)
    ws["A3"] = "Project: Ashburn DC-14 Electrical Rough-In (200k SF)"
    ws["A4"] = "Contract Value: $4,250,000"
    ws["A5"] = "Duration: 14 months"
    ws["A6"] = "Discount Rate: 10% annual"
    ws["A8"] = "Tabs in this workbook:"
    for i, t in enumerate(["Assumptions", "Cash Flow", "Sensitivity", "Summary", "Executive"], start=9):
        ws[f"A{i}"] = f"  • {t}"

    # Assumptions
    asm = wb.create_sheet("Assumptions")
    asm.append(["Item", "Value", "Notes"])
    style_header_row(asm, 1, 3)
    rows = [
        ("Contract Value", 4250000, "GMP contract"),
        ("Labor %", 0.42, "Journeymen + apprentices"),
        ("Materials %", 0.31, "Switchgear, conduit, wire"),
        ("Subs %", 0.12, "Testing, commissioning"),
        ("Overhead %", 0.08, "PM, admin, truck"),
        ("Target Margin", 0.07, "After all costs"),
        ("Discount Rate (annual)", 0.10, "WACC"),
        ("Duration (months)", 14, ""),
    ]
    for r in rows:
        asm.append(r)
    auto_width(asm, ["Item", "Value", "Notes"])

    # Cash flow
    cf = wb.create_sheet("Cash Flow")
    headers = ["Month", "Billing", "Labor", "Materials", "Subs", "OH", "Net Cash", "Cumulative"]
    cf.append(headers)
    style_header_row(cf, 1, len(headers))
    cum = 0
    months = 14
    for m in range(1, months + 1):
        bill = 4250000 / months
        lab = bill * 0.42
        mat = bill * 0.31
        sub = bill * 0.12
        oh = bill * 0.08
        net = bill - lab - mat - sub - oh
        cum += net
        cf.append([m, round(bill), round(lab), round(mat), round(sub), round(oh), round(net), round(cum)])
    auto_width(cf, headers)

    # Sensitivity
    sens = wb.create_sheet("Sensitivity")
    sens.append(["Cost Overrun Scenario", "NPV", "IRR", "Go/No-Go"])
    style_header_row(sens, 1, 4)
    scenarios = [
        ("Baseline (0%)", 287500, 0.184, "GO"),
        ("+10% overrun", 157200, 0.124, "GO"),
        ("+20% overrun", 26900, 0.068, "MARGINAL"),
        ("+30% overrun", -103400, 0.012, "NO-GO"),
    ]
    for s in scenarios:
        sens.append(s)
    auto_width(sens, ["Scenario", "NPV", "IRR", "Go"])

    # Summary
    summ = wb.create_sheet("Summary")
    summ["A1"] = "Recommendation: GO"
    summ["A1"].font = Font(bold=True, size=14, color="059669")
    summ["A3"] = "Baseline NPV: $287,500"
    summ["A4"] = "Baseline IRR: 18.4%"
    summ["A5"] = "Payback Period: 9.2 months"
    summ["A6"] = "Break-even overrun: +22%"

    ex = wb.create_sheet("Executive Summary")
    ex["A1"] = "Executive Summary"
    ex["A1"].font = Font(bold=True, size=16)
    ex["A3"] = (
        "The Ashburn DC-14 electrical rough-in generates $287,500 in risk-adjusted value "
        "at the 10% discount rate, with an 18.4% IRR. The project remains profitable through "
        "+20% cost overruns; above +22% overrun, NPV turns negative."
    )
    ex["A3"].alignment = Alignment(wrap_text=True, vertical="top")
    ex.column_dimensions["A"].width = 90

    wb.save(OUT / "Patriot_Electric_NPV.xlsx")


# --- 2. Budget -------------------------------------------------------------
def build_budget():
    wb = Workbook()
    ws = wb.active
    ws.title = "Monthly Budget"
    headers = ["Category", "Monthly Target", "Q1 Actual", "Variance", "% of Revenue"]
    ws.append(headers)
    style_header_row(ws, 1, len(headers))
    rows = [
        ("Revenue", 412000, 438200, 26200, 1.00),
        ("Labor (W-2)", 173000, 181400, -8400, 0.414),
        ("Materials", 128000, 134700, -6700, 0.307),
        ("Subcontractors", 49000, 52100, -3100, 0.119),
        ("Insurance & Bonds", 18500, 18500, 0, 0.042),
        ("Vehicles & Fuel", 11200, 12800, -1600, 0.029),
        ("Office & Admin", 9400, 9200, 200, 0.021),
        ("Owner Comp", 16000, 16000, 0, 0.037),
        ("Net Profit", 6900, 13500, 6600, 0.031),
    ]
    for r in rows:
        ws.append(r)
    auto_width(ws, headers)
    wb.save(OUT / "Patriot_Electric_Budget.xlsx")


# --- 3. Financial Model ----------------------------------------------------
def build_finmodel():
    wb = Workbook()
    pl = wb.active
    pl.title = "P&L"
    pl.append(["Line Item", "2024", "2025", "2026E"])
    style_header_row(pl, 1, 4)
    lines = [
        ("Revenue", 3800000, 4650000, 5400000),
        ("COGS", 2890000, 3488000, 4050000),
        ("Gross Profit", 910000, 1162000, 1350000),
        ("SG&A", 684000, 802000, 918000),
        ("EBITDA", 226000, 360000, 432000),
        ("Net Income", 182000, 301000, 368000),
    ]
    for l in lines:
        pl.append(l)
    auto_width(pl, ["Item", "2024", "2025", "2026"])

    kpi = wb.create_sheet("KPI Dashboard")
    kpi.append(["Metric", "Value"])
    style_header_row(kpi, 1, 2)
    kpis = [
        ("Gross Margin", "25.0%"),
        ("EBITDA Margin", "8.0%"),
        ("Revenue Growth YoY", "+16.1%"),
        ("Avg Job Size", "$84,375"),
        ("Job Close Rate", "34%"),
        ("Days Sales Outstanding", "48"),
    ]
    for k in kpis:
        kpi.append(k)
    auto_width(kpi, ["Metric", "Value"])
    wb.save(OUT / "Patriot_Electric_FinancialModel.xlsx")


# --- 4. Project Estimate ---------------------------------------------------
def build_estimate():
    wb = Workbook()
    ws = wb.active
    ws.title = "Estimate"
    headers = ["Item", "Qty", "Unit", "Unit Cost", "Extended"]
    ws.append(headers)
    style_header_row(ws, 1, len(headers))
    items = [
        ("200A main switchgear", 1, "EA", 18400, 18400),
        ("75kVA step-down transformer", 3, "EA", 4850, 14550),
        ("Branch panelboards 42-circuit", 8, "EA", 2180, 17440),
        ("600A feeder, 200ft", 1, "LOT", 11200, 11200),
        ("Conduit + wire rough-in", 1, "LOT", 38400, 38400),
        ("Lighting rough-in (LED 2x4)", 160, "EA", 142, 22720),
        ("Device rough-in", 340, "EA", 68, 23120),
        ("Testing & commissioning", 1, "LOT", 8600, 8600),
        ("Labor (journeyman + helper)", 840, "HR", 98, 82320),
        ("Project management", 1, "LOT", 14200, 14200),
    ]
    for it in items:
        ws.append(it)
    total = sum(i[4] for i in items)
    ws.append(["", "", "", "Subtotal", total])
    ws.append(["", "", "", "Overhead 8%", round(total * 0.08)])
    ws.append(["", "", "", "Margin 12%", round(total * 0.12)])
    ws.append(["", "", "", "Total", round(total * 1.20)])
    auto_width(ws, headers)
    wb.save(OUT / "Patriot_Electric_Estimate.xlsx")


# --- 5. Proposal (DOCX) ----------------------------------------------------
def build_proposal():
    doc = Document()
    t = doc.add_heading("Proposal - Ashburn DC-14 Electrical Rough-In", level=1)
    doc.add_paragraph("Prepared for: Meridian Construction, Inc.")
    doc.add_paragraph("From: Patriot Electric, LLC  ·  Ashburn, VA  ·  License #2705-118432")
    doc.add_paragraph("Date: April 5, 2026")

    doc.add_heading("Scope of Work", level=2)
    doc.add_paragraph(
        "Patriot Electric will furnish all labor, materials, tools, and supervision necessary to "
        "complete the electrical rough-in for the Ashburn DC-14 data center shell, including "
        "main switchgear installation, step-down transformers, branch panelboards, feeder runs, "
        "lighting rough-in, device rough-in, and commissioning."
    )

    doc.add_heading("Timeline", level=2)
    doc.add_paragraph("Mobilization: May 1, 2026")
    doc.add_paragraph("Rough-in complete: September 15, 2026")
    doc.add_paragraph("Commissioning: October 1, 2026")

    doc.add_heading("Pricing", level=2)
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Light Grid Accent 1"
    tbl.rows[0].cells[0].text = "Line Item"
    tbl.rows[0].cells[1].text = "Amount"
    for label, amt in [
        ("Materials & Equipment", "$135,830"),
        ("Labor", "$82,320"),
        ("PM + Commissioning", "$22,800"),
        ("Overhead", "$19,280"),
        ("Margin", "$28,920"),
        ("Total", "$289,150"),
    ]:
        row = tbl.add_row().cells
        row[0].text = label
        row[1].text = amt

    doc.add_heading("Terms", level=2)
    doc.add_paragraph("Net 30 from invoice date. 10% mobilization, monthly progress billing.")

    doc.save(OUT / "Patriot_Electric_Proposal.docx")


# --- 6. Capabilities Doc (PPTX) --------------------------------------------
def build_capabilities():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Patriot Electric"
    slide.placeholders[1].text = "Ashburn, VA  ·  Data Center Electrical Specialists"

    for title, bullets in [
        ("Who We Are", [
            "Licensed electrical contractor based in Ashburn, VA",
            "Focused on data center corridor (Loudoun, Prince William)",
            "22 journeymen, 8 apprentices, 4 project managers",
            "$4.6M revenue (2025), targeting $5.4M (2026)",
        ]),
        ("Recent Projects", [
            "DC-11 Electrical Rough-In - $2.1M (2025)",
            "DC-09 Branch Panel Upgrade - $685K (2025)",
            "Loudoun Logistics Park - $1.4M (2024)",
        ]),
        ("Capabilities", [
            "Medium-voltage switchgear up to 15kV",
            "Generator + UPS integration",
            "Structured cable pathways",
            "NFPA 70E arc-flash studies",
        ]),
        ("Why Patriot", [
            "Data-center-only focus - we know the playbook",
            "100% on-time completion, 2024 + 2025",
            "Zero OSHA recordables, 3.8M safe hours",
        ]),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b

    prs.save(OUT / "Patriot_Electric_Capabilities.pptx")


# --- 7. Competitive Analysis (PDF) -----------------------------------------
def build_competitive():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "Competitive Analysis - Ashburn DC Corridor", ln=1)
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 6, "Prepared for Patriot Electric  |  April 2026", ln=1)
    pdf.ln(6)

    competitors = [
        ("Power Design Inc.",
         "National MEP. Strong on mega-projects >$50M. Weak on sub-$5M scopes - slow to respond, high overhead.",
         "Our edge: faster mobilization, sub-$5M focus, local relationships."),
        ("Mona Electric Group",
         "Regional leader. Strong bench, good pricing, strong Meridian relationship.",
         "Our edge: tighter schedule commitments, owner-present on site."),
        ("Rosendin Electric",
         "Tier-1 national. Premium pricing, best-in-class safety record, strong labor pipeline.",
         "Our edge: 15-20% lower cost on sub-$3M rough-ins."),
    ]
    for name, profile, edge in competitors:
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, name, ln=1)
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(190, 5, f"Profile: {profile}")
        pdf.multi_cell(190, 5, f"Our edge: {edge}")
        pdf.ln(3)

    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, "Positioning Recommendation", ln=1)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(190, 5,
        "Patriot Electric should position as the 'responsive sub-$5M DC specialist.' "
        "Emphasize schedule certainty, owner-on-site, and local crew continuity in all "
        "Meridian and Compass bid responses.")

    pdf.output(str(OUT / "Patriot_Electric_CompetitiveAnalysis.pdf"))


# --- 8. Financial Review (PDF) ---------------------------------------------
def build_review():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "Monthly Financial Review - March 2026", ln=1)
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 6, "Patriot Electric, LLC", ln=1)
    pdf.ln(6)

    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, "Headline", ln=1)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(190, 5,
        "March revenue $438K, +6.4% over target. Net margin 3.1% vs 1.7% target. "
        "Labor ran hot ($8.4K over) on the DC-11 punch list.")
    pdf.ln(3)

    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, "Action Items", ln=1)
    pdf.set_font("Helvetica", "", 10)
    for i, a in enumerate([
        "Close DC-11 punch list by April 12 - $18K labor at risk",
        "Invoice Meridian for DC-09 retention ($42K) - 38 days outstanding",
        "Re-bid DC-14 with +2% contingency before April 20",
        "Schedule Q2 safety refresher for new apprentices",
    ], start=1):
        pdf.cell(0, 6, f"{i}. {a}", ln=1)

    pdf.output(str(OUT / "Patriot_Electric_FinancialReview.pdf"))


if __name__ == "__main__":
    build_npv()
    build_budget()
    build_finmodel()
    build_estimate()
    build_proposal()
    build_capabilities()
    build_competitive()
    build_review()
    print(f"Wrote demo files to {OUT}")
    for f in sorted(OUT.iterdir()):
        print(f"  {f.name}  ({f.stat().st_size:,} bytes)")
